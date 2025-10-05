"""
Text Extraction Module
Extracts text from various file formats including images, PDFs, Office documents
"""
import os
from typing import Dict, Optional

try:
    from PIL import Image
    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False

try:
    import pytesseract
    PYTESSERACT_AVAILABLE = True
except ImportError:
    PYTESSERACT_AVAILABLE = False

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class TextExtractor:
    """Extracts text from multiple file formats"""
    
    def __init__(self):
        self.supported_formats = {
            'image': ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif'],
            'pdf': ['.pdf'],
            'word': ['.docx', '.doc'],
            'excel': ['.xlsx', '.xls'],
            'powerpoint': ['.pptx', '.ppt'],
            'text': ['.txt']
        }
    
    def extract_from_image(self, file_path: str) -> str:
        """
        Extract text from image using OCR
        
        Args:
            file_path: Path to image file
            
        Returns:
            Extracted text
        """
        if not PILLOW_AVAILABLE or not PYTESSERACT_AVAILABLE:
            return "Error: PIL or pytesseract not installed. Install with: pip install Pillow pytesseract"
        
        try:
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            return text.strip()
        except Exception as e:
            return f"Error extracting text from image: {str(e)}"
    
    def extract_from_pdf(self, file_path: str) -> str:
        """
        Extract text from PDF file
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Extracted text
        """
        if not PYPDF2_AVAILABLE:
            return "Error: PyPDF2 not installed. Install with: pip install PyPDF2"
        
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            return f"Error extracting text from PDF: {str(e)}"
    
    def extract_from_word(self, file_path: str) -> str:
        """
        Extract text from Word document
        
        Args:
            file_path: Path to Word file
            
        Returns:
            Extracted text
        """
        if not DOCX_AVAILABLE:
            return "Error: python-docx not installed. Install with: pip install python-docx"
        
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            return f"Error extracting text from Word: {str(e)}"
    
    def extract_from_excel(self, file_path: str) -> str:
        """
        Extract text from Excel spreadsheet
        
        Args:
            file_path: Path to Excel file
            
        Returns:
            Extracted text
        """
        if not OPENPYXL_AVAILABLE:
            return "Error: openpyxl not installed. Install with: pip install openpyxl"
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n=== Sheet: {sheet_name} ===\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
            
            return text.strip()
        except Exception as e:
            return f"Error extracting text from Excel: {str(e)}"
    
    def extract_from_powerpoint(self, file_path: str) -> str:
        """
        Extract text from PowerPoint presentation
        
        Args:
            file_path: Path to PowerPoint file
            
        Returns:
            Extracted text
        """
        if not PPTX_AVAILABLE:
            return "Error: python-pptx not installed. Install with: pip install python-pptx"
        
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n=== Slide {slide_num} ===\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return text.strip()
        except Exception as e:
            return f"Error extracting text from PowerPoint: {str(e)}"
    
    def extract_from_text(self, file_path: str) -> str:
        """
        Extract text from plain text file
        
        Args:
            file_path: Path to text file
            
        Returns:
            Extracted text
        """
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            return text.strip()
        except Exception as e:
            return f"Error extracting text from text file: {str(e)}"
    
    def extract_text(self, file_path: str) -> Dict[str, any]:
        """
        Extract text from any supported file format
        
        Args:
            file_path: Path to file
            
        Returns:
            Dictionary with extracted text and metadata
        """
        _, ext = os.path.splitext(file_path.lower())
        
        result = {
            'file_path': file_path,
            'file_type': ext,
            'text': '',
            'success': False
        }
        
        try:
            if ext in self.supported_formats['image']:
                result['text'] = self.extract_from_image(file_path)
                result['success'] = True
            elif ext in self.supported_formats['pdf']:
                result['text'] = self.extract_from_pdf(file_path)
                result['success'] = True
            elif ext in self.supported_formats['word']:
                result['text'] = self.extract_from_word(file_path)
                result['success'] = True
            elif ext in self.supported_formats['excel']:
                result['text'] = self.extract_from_excel(file_path)
                result['success'] = True
            elif ext in self.supported_formats['powerpoint']:
                result['text'] = self.extract_from_powerpoint(file_path)
                result['success'] = True
            elif ext in self.supported_formats['text']:
                result['text'] = self.extract_from_text(file_path)
                result['success'] = True
            else:
                result['text'] = f"Unsupported file format: {ext}"
                result['success'] = False
        except Exception as e:
            result['text'] = f"Error processing file: {str(e)}"
            result['success'] = False
        
        return result
